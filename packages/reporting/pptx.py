"""PowerPoint generation (Milestone 5) via python-pptx: a 1-slide
executive summary and a 2-slide ELT board pack. Same deferred-scope note
as `pdf.py` — always the full current-state dashboard, `scope`/`period`
are cover-page metadata only.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from packages.reporting.data import ReportContext

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

INK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x64, 0x74, 0x8B)
BAND_COLORS = {
    "low": RGBColor(0x1A, 0x7F, 0x37),
    "moderate": RGBColor(0x9A, 0x67, 0x00),
    "high": RGBColor(0xBC, 0x4C, 0x00),
    "extreme": RGBColor(0xCF, 0x22, 0x2E),
}

BLANK_LAYOUT_INDEX = 6


def _new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def _add_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])


def _add_title(slide, text: str, *, top=Inches(0.35)):
    box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(0.7))
    tf = box.text_frame
    tf.text = text
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = INK
    return box


def _add_subtitle(slide, text: str, *, top=Inches(1.0)):
    box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(0.4))
    tf = box.text_frame
    tf.text = text
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(13)
    run.font.color.rgb = MUTED
    return box


def _meta_line(context: ReportContext) -> str:
    meta = f"Generated {context.generated_at.strftime('%Y-%m-%d %H:%M UTC')} · Scope: {context.scope_label}"
    if context.period_start or context.period_end:
        meta += f" · Period: {context.period_start or '—'} to {context.period_end or '—'}"
    return meta


def _add_kpi_row(slide, kpis: list[tuple[str, str | int]], *, top, colors=None):
    n = len(kpis)
    tile_width = Inches(12.3) / n
    for i, (label, value) in enumerate(kpis):
        left = Inches(0.5) + tile_width * i
        box = slide.shapes.add_textbox(left, top, tile_width - Inches(0.15), Inches(1.1))
        tf = box.text_frame
        tf.word_wrap = True
        p_value = tf.paragraphs[0]
        p_value.text = str(value)
        p_value.runs[0].font.size = Pt(32)
        p_value.runs[0].font.bold = True
        p_value.runs[0].font.color.rgb = (colors or {}).get(label, INK)
        p_label = tf.add_paragraph()
        p_label.text = label
        p_label.runs[0].font.size = Pt(11)
        p_label.runs[0].font.color.rgb = MUTED


def _add_table(slide, rows: list[list[str]], *, left, top, width, height, band_col: int | None = None):
    n_rows, n_cols = len(rows), len(rows[0])
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table
    for c, header in enumerate(rows[0]):
        cell = table.cell(0, c)
        cell.text = header
        cell.text_frame.paragraphs[0].runs[0].font.bold = True
        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
    for r in range(1, n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.text = rows[r][c]
            runs = cell.text_frame.paragraphs[0].runs
            if not runs:
                continue  # an empty cell value produces a paragraph with no run
            run = runs[0]
            run.font.size = Pt(10)
            if band_col is not None and c == band_col:
                color = BAND_COLORS.get(rows[r][c].lower())
                if color:
                    run.font.color.rgb = color
                    run.font.bold = True
    return table


def render_pptx_one_slide(context: ReportContext, output_path: str | Path) -> None:
    prs = _new_presentation()
    slide = _add_slide(prs)
    d = context.dashboard

    _add_title(slide, "Risk Intelligence Platform — Executive Summary")
    _add_subtitle(slide, _meta_line(context))

    _add_kpi_row(
        slide,
        [
            ("Total risks", d["total_risks"]),
            ("Extreme", d["extreme_count"]),
            ("High", d["high_count"]),
            ("Weak controls", d["weak_controls_count"]),
            ("Overdue actions", d["overdue_actions_count"]),
            ("Outside appetite", d["risks_outside_appetite_count"]),
        ],
        top=Inches(1.6),
        colors={"Extreme": BAND_COLORS["extreme"], "High": BAND_COLORS["high"]},
    )

    top_rows = [["Risk", "Band", "Owner"]]
    for r in d["top_risks"][:6]:
        top_rows.append([r["title"], (r["residual_band"] or "—").capitalize(), r["owner_email"] or "Unassigned"])
    if len(top_rows) == 1:
        top_rows.append(["No scored risks yet.", "", ""])
    _add_table(
        slide, top_rows, left=Inches(0.5), top=Inches(3.1), width=Inches(12.3), height=Inches(3.5), band_col=1
    )

    prs.save(str(output_path))


def render_pptx_two_slide_elt(context: ReportContext, output_path: str | Path) -> None:
    prs = _new_presentation()
    d = context.dashboard

    slide1 = _add_slide(prs)
    _add_title(slide1, "Risk Overview")
    _add_subtitle(slide1, _meta_line(context))
    _add_kpi_row(
        slide1,
        [
            ("Total risks", d["total_risks"]),
            ("Extreme", d["extreme_count"]),
            ("High", d["high_count"]),
            ("Moderate", d["moderate_count"]),
            ("Low", d["low_count"]),
        ],
        top=Inches(1.6),
        colors={"Extreme": BAND_COLORS["extreme"], "High": BAND_COLORS["high"]},
    )
    cat_rows = [["Category", "Risk count", "Avg residual score"]]
    for c in d["category_exposure"][:8]:
        cat_rows.append(
            [
                c["category_name"],
                str(c["risk_count"]),
                str(c["avg_residual_score"]) if c["avg_residual_score"] is not None else "—",
            ]
        )
    if len(cat_rows) == 1:
        cat_rows.append(["No categorized risks.", "", ""])
    _add_table(slide1, cat_rows, left=Inches(0.5), top=Inches(3.1), width=Inches(12.3), height=Inches(3.5))

    slide2 = _add_slide(prs)
    _add_title(slide2, "Governance & Leadership Attention")
    _add_kpi_row(
        slide2,
        [
            ("Weak controls", d["weak_controls_count"]),
            ("Overdue actions", d["overdue_actions_count"]),
            ("Overdue reviews", d["overdue_reviews_count"]),
            ("Outside appetite", d["risks_outside_appetite_count"]),
        ],
        top=Inches(1.3),
    )
    top_rows = [["Risk", "Category", "Residual", "Band", "Owner"]]
    for r in d["top_risks"][:8]:
        top_rows.append(
            [
                r["title"],
                r["category_name"] or "Uncategorized",
                str(r["residual_score"]) if r["residual_score"] is not None else "—",
                (r["residual_band"] or "—").capitalize(),
                r["owner_email"] or "Unassigned",
            ]
        )
    if len(top_rows) == 1:
        top_rows.append(["No scored risks yet.", "", "", "", ""])
    _add_table(
        slide2, top_rows, left=Inches(0.5), top=Inches(2.6), width=Inches(12.3), height=Inches(4), band_col=3
    )

    prs.save(str(output_path))
