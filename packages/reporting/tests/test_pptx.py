from datetime import date, datetime, timezone

from pptx import Presentation

from packages.reporting.data import ReportContext
from packages.reporting.pptx import render_pptx_one_slide, render_pptx_two_slide_elt

SAMPLE_DASHBOARD = {
    "total_risks": 3,
    "extreme_count": 1,
    "high_count": 1,
    "moderate_count": 1,
    "low_count": 0,
    "unscored_count": 0,
    "overdue_reviews_count": 2,
    "weak_controls_count": 1,
    "overdue_actions_count": 4,
    "risks_outside_appetite_count": 1,
    "top_risks": [
        {
            "id": "r1",
            "risk_code": "RSK-0001",
            "title": "Extreme sample risk",
            "category_name": "Cyber & Information Security",
            "residual_score": 12.5,
            "residual_band": "extreme",
            "owner_email": "owner@example.com",
            "next_review_date": "2026-01-01",
        },
    ],
    "category_exposure": [
        {"category_id": "c1", "category_name": "Cyber & Information Security", "risk_count": 3, "avg_residual_score": 8.1},
    ],
}


def _context(dashboard=None) -> ReportContext:
    return ReportContext(
        generated_at=datetime(2026, 6, 1, 12, 0, tzinfo=timezone.utc),
        period_start=None,
        period_end=None,
        scope_label="All risks",
        dashboard=dashboard if dashboard is not None else SAMPLE_DASHBOARD,
    )


def _all_text(slide) -> str:
    return " ".join(
        shape.text_frame.text
        for shape in slide.shapes
        if shape.has_text_frame and shape.text_frame.text
    )


class TestRenderPptxOneSlide:
    def test_writes_one_slide_with_kpis_and_top_risks_table(self, tmp_path):
        output_path = tmp_path / "one_slide.pptx"
        render_pptx_one_slide(_context(), output_path)

        prs = Presentation(output_path)
        assert len(prs.slides) == 1

        slide = prs.slides[0]
        assert "Executive Summary" in _all_text(slide)

        tables = [shape.table for shape in slide.shapes if shape.has_table]
        assert len(tables) == 1
        rows = [[cell.text for cell in row.cells] for row in tables[0].rows]
        assert rows[0] == ["Risk", "Band", "Owner"]
        assert rows[1][0] == "Extreme sample risk"

    def test_handles_no_scored_risks(self, tmp_path):
        empty_dashboard = {**SAMPLE_DASHBOARD, "top_risks": []}
        output_path = tmp_path / "empty_one_slide.pptx"

        render_pptx_one_slide(_context(empty_dashboard), output_path)

        prs = Presentation(output_path)
        tables = [shape.table for shape in prs.slides[0].shapes if shape.has_table]
        rows = [[cell.text for cell in row.cells] for row in tables[0].rows]
        assert rows[1][0] == "No scored risks yet."


class TestRenderPptxTwoSlideElt:
    def test_writes_two_slides(self, tmp_path):
        output_path = tmp_path / "two_slide.pptx"
        render_pptx_two_slide_elt(_context(), output_path)

        prs = Presentation(output_path)
        assert len(prs.slides) == 2

        slide1_text = _all_text(prs.slides[0])
        assert "Risk Overview" in slide1_text

        slide2_text = _all_text(prs.slides[1])
        assert "Governance" in slide2_text

        slide1_tables = [shape.table for shape in prs.slides[0].shapes if shape.has_table]
        cat_rows = [[cell.text for cell in row.cells] for row in slide1_tables[0].rows]
        assert cat_rows[1][0] == "Cyber & Information Security"

        slide2_tables = [shape.table for shape in prs.slides[1].shapes if shape.has_table]
        risk_rows = [[cell.text for cell in row.cells] for row in slide2_tables[0].rows]
        assert risk_rows[1][0] == "Extreme sample risk"
